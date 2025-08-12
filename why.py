# CoorSight Document Analysis Platform - Enhanced Multi-Format Support with Advanced Contradiction Detection
import streamlit.components.v1 as components
import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime, timedelta
import uuid
import re
import io
import json
import csv
from typing import Dict, List, Any
import plotly.express as px

# Enhanced imports for contradiction detection
try:
    import spacy
    nlp = spacy.load("en_core_web_sm")
    SPACY_AVAILABLE = True
except:
    SPACY_AVAILABLE = False

try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMER_AVAILABLE = True
except:
    SENTENCE_TRANSFORMER_AVAILABLE = False

# Document processing imports
try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except:
    XLSX_AVAILABLE = False

try:
    import xlrd
    XLS_AVAILABLE = True
except:
    XLS_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except:
    RTF_AVAILABLE = False

try:
    import xml.etree.ElementTree as ET
    XML_AVAILABLE = True
except:
    XML_AVAILABLE = False

# =============================================================================
# PAGE CONFIGURATION
# =============================================================================

st.set_page_config(
    page_title="CoorSight Document Analysis",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# =============================================================================
# SESSION STATE INITIALIZATION
# =============================================================================

if 'documents' not in st.session_state:
    st.session_state.documents = {}
if 'analyses' not in st.session_state:
    st.session_state.analyses = {}
if 'current_page' not in st.session_state:
    st.session_state.current_page = "📤 Upload Documents"

# =============================================================================
# FILE PROCESSING FUNCTIONS - ENHANCED (Keep your existing functions)
# =============================================================================

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF files using PyMuPDF"""
    if not PDF_AVAILABLE:
        return "PDF processing not available. Install PyMuPDF: pip install PyMuPDF"
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        text = "".join([page.get_text() for page in doc])
        doc.close()
        return text.strip()
    except Exception as e:
        return f"Error processing PDF: {str(e)}"

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX files using python-docx"""
    if not DOCX_AVAILABLE:
        return "DOCX processing not available. Install python-docx: pip install python-docx"
    try:
        doc = Document(io.BytesIO(file_content))
        text_parts = []
        
        # Extract paragraph text
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        
        return "\n".join(text_parts)
    except Exception as e:
        return f"Error processing DOCX: {str(e)}"

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint files using python-pptx"""
    if not PPTX_AVAILABLE:
        return "PPTX processing not available. Install python-pptx: pip install python-pptx"
    try:
        prs = Presentation(io.BytesIO(file_content))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {slide_num} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                
                # Extract text from tables in slides
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
        
        return "\n".join(text_parts)
    except Exception as e:
        return f"Error processing PPTX: {str(e)}"

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from Excel XLSX files using openpyxl"""
    if not XLSX_AVAILABLE:
        return "XLSX processing not available. Install openpyxl: pip install openpyxl"
    try:
        workbook = load_workbook(io.BytesIO(file_content), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n--- Sheet: {sheet_name} ---")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None and str(cell).strip():
                        row_text.append(str(cell))
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n".join(text_parts)
    except Exception as e:
        return f"Error processing XLSX: {str(e)}"

def extract_text_from_xls(file_content: bytes) -> str:
    """Extract text from Excel XLS files using xlrd"""
    if not XLS_AVAILABLE:
        return "XLS processing not available. Install xlrd: pip install xlrd"
    try:
        workbook = xlrd.open_workbook(file_contents=file_content)
        text_parts = []
        
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            text_parts.append(f"\n--- Sheet: {sheet.name} ---")
            
            for row_idx in range(sheet.nrows):
                row_text = []
                for col_idx in range(sheet.ncols):
                    cell = sheet.cell(row_idx, col_idx)
                    if cell.value and str(cell.value).strip():
                        row_text.append(str(cell.value))
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n".join(text_parts)
    except Exception as e:
        return f"Error processing XLS: {str(e)}"

def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from CSV files"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            return "Error: Could not decode CSV file with any supported encoding"
        
        # Parse CSV
        csv_reader = csv.reader(io.StringIO(text_content))
        text_parts = []
        
        for row_idx, row in enumerate(csv_reader):
            if row:  # Skip empty rows
                if row_idx == 0:  # Header row
                    text_parts.append(f"Headers: {' | '.join(row)}")
                else:
                    text_parts.append(" | ".join(str(cell) for cell in row if str(cell).strip()))
        
        return "\n".join(text_parts)
    except Exception as e:
        return f"Error processing CSV: {str(e)}"

def extract_text_from_json(file_content: bytes) -> str:
    """Extract text from JSON files"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            return "Error: Could not decode JSON file"
        
        # Parse JSON and extract text
        data = json.loads(text_content)
        
        def extract_text_from_json_recursive(obj, level=0):
            text_parts = []
            indent = "  " * level
            
            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, (dict, list)):
                        text_parts.append(f"{indent}{key}:")
                        text_parts.extend(extract_text_from_json_recursive(value, level + 1))
                    else:
                        text_parts.append(f"{indent}{key}: {value}")
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    if isinstance(item, (dict, list)):
                        text_parts.append(f"{indent}[{i}]:")
                        text_parts.extend(extract_text_from_json_recursive(item, level + 1))
                    else:
                        text_parts.append(f"{indent}[{i}]: {item}")
            else:
                text_parts.append(f"{indent}{obj}")
            
            return text_parts
        
        return "\n".join(extract_text_from_json_recursive(data))
    except Exception as e:
        return f"Error processing JSON: {str(e)}"

def extract_text_from_rtf(file_content: bytes) -> str:
    """Extract text from RTF files"""
    if not RTF_AVAILABLE:
        return "RTF processing not available. Install striprtf: pip install striprtf"
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                rtf_content = file_content.decode(encoding)
                return rtf_to_text(rtf_content)
            except UnicodeDecodeError:
                continue
        
        return "Error: Could not decode RTF file"
    except Exception as e:
        return f"Error processing RTF: {str(e)}"

def extract_text_from_xml(file_content: bytes) -> str:
    """Extract text from XML files"""
    if not XML_AVAILABLE:
        return "XML processing not available"
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                xml_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            return "Error: Could not decode XML file"
        
        root = ET.fromstring(xml_content)
        
        def extract_xml_text(element, level=0):
            text_parts = []
            indent = "  " * level
            
            # Add element tag and text
            if element.text and element.text.strip():
                text_parts.append(f"{indent}{element.tag}: {element.text.strip()}")
            else:
                text_parts.append(f"{indent}{element.tag}")
            
            # Add attributes
            if element.attrib:
                for key, value in element.attrib.items():
                    text_parts.append(f"{indent}  @{key}: {value}")
            
            # Process child elements
            for child in element:
                text_parts.extend(extract_xml_text(child, level + 1))
            
            return text_parts
        
        return "\n".join(extract_xml_text(root))
    except Exception as e:
        return f"Error processing XML: {str(e)}"

def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from TXT files with multiple encoding support"""
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            return file_content.decode(encoding)
        except UnicodeDecodeError:
            continue
    
    return "Error: Could not decode text file with any supported encoding"

def process_file(uploaded_file) -> Dict[str, Any]:
    """Process uploaded file and extract text based on file type"""
    file_content = uploaded_file.read()
    file_ext = uploaded_file.name.split('.')[-1].lower()
    
    # Define file type processors
    processors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_docx,  # Treat .doc same as .docx
        'pptx': extract_text_from_pptx,
        'ppt': extract_text_from_pptx,  # Treat .ppt same as .pptx
        'xlsx': extract_text_from_xlsx,
        'xls': extract_text_from_xls,
        'csv': extract_text_from_csv,
        'json': extract_text_from_json,
        'rtf': extract_text_from_rtf,
        'xml': extract_text_from_xml,
        'txt': extract_text_from_txt,
        'md': extract_text_from_txt,
        'log': extract_text_from_txt,
    }
    
    # Process file
    if file_ext in processors:
        text = processors[file_ext](file_content)
    else:
        text = f"File type '{file_ext}' not supported. Supported types: {', '.join(processors.keys())}"
    
    # Generate document metadata
    doc_id = str(uuid.uuid4())
    return {
        'id': doc_id,
        'name': uploaded_file.name,
        'type': file_ext,
        'size': len(file_content),
        'text': text,
        'uploaded_at': datetime.now(),
        'word_count': len(text.split()) if text and not text.startswith('Error') else 0,
        'has_error': text.startswith('Error') if text else True
    }

# =============================================================================
# ENHANCED CONTRADICTION DETECTION INTEGRATION
# =============================================================================

# Import the enhanced contradiction detection
def find_contradictions_enhanced(documents: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Enhanced contradiction detection function that integrates with the enhanced detection module.
    This function will import the enhanced detection when called.
    """
    try:
        # Try to import the enhanced detection module
        from enhanced_contradiction_detection import find_contradictions_enhanced as enhanced_detect
        return enhanced_detect(documents)
    except ImportError:
        # Fallback to basic detection if enhanced module not available
        st.warning("⚠️ Enhanced contradiction detection module not found. Using basic detection.")
        return find_contradictions_basic(documents)

def find_contradictions_basic(documents: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Fallback basic contradiction detection"""
    contradictions = []
    patterns = [
        (r'deadline is (\w+)', 'deadline'),
        (r'budget is \$?([\d,]+)', 'budget'),
        (r'completed in (\d+) days', 'timeline')
    ]
    
    texts = [doc['text'] for doc in documents if not doc.get('has_error', False)]
    
    for i, text1 in enumerate(texts):
        for j, text2 in enumerate(texts[i+1:], i+1):
            for pattern, topic in patterns:
                matches1 = re.findall(pattern, text1.lower())
                matches2 = re.findall(pattern, text2.lower())
                if matches1 and matches2 and matches1[0] != matches2[0]:
                    contradictions.append({
                        'type': topic,
                        'description': f'{topic.title()} contradiction detected',
                        'doc1_name': documents[i]['name'],
                        'doc1_value': matches1[0],
                        'doc2_name': documents[j]['name'],
                        'doc2_value': matches2[0],
                        'severity': 'Medium',
                        'confidence': 0.7,
                        'suggestion': f'Review {topic} values for accuracy',
                        'entities': [],
                        'context': {'type': 'basic_pattern'}
                    })
    
    return contradictions

# =============================================================================
# ENHANCED CONTRADICTION DETECTION PAGE
# =============================================================================

def page_contradiction_detection():
    """Enhanced Contradiction Detection Page"""
    st.header("⚠️ Enhanced Contradiction Detection")
    st.write("Find contradictions and inconsistencies across multiple documents using advanced AI techniques")

    # Show AI enhancement status
    col1, col2, col3 = st.columns(3)
    with col1:
        if SPACY_AVAILABLE:
            st.success("✅ spaCy NLP Available")
        else:
            st.warning("❌ spaCy NLP Missing")
    
    with col2:
        if SENTENCE_TRANSFORMER_AVAILABLE:
            st.success("✅ Semantic Analysis Available")
        else:
            st.warning("❌ Semantic Analysis Missing")
    
    with col3:
        try:
            from enhanced_contradiction_detection import EnhancedContradictionDetector
            st.success("✅ Enhanced Detection Ready")
        except ImportError:
            st.error("❌ Enhanced Module Missing")

    if len(st.session_state.documents) < 2:
        st.info("Need at least 2 documents to detect contradictions.")
        st.markdown("""
        ### 🎯 What This Enhanced Tool Detects:
        - **Temporal Contradictions** ⏰: Conflicting dates, deadlines, and timelines
        - **Numerical Contradictions** 🔢: Different budgets, quantities, and measurements  
        - **Factual Contradictions** 📋: Conflicting statements about the same entities
        - **Status Contradictions** 📊: Different completion or approval statuses
        - **Semantic Contradictions** 🧠: Meaning conflicts detected using AI
        
        ### 📦 For Full AI Features, Install:
        ```bash
        pip install spacy sentence-transformers
        python -m spacy download en_core_web_sm
        ```
        """)
        return

    file_type_icons = {
        'pdf': '📄', 'docx': '📝', 'doc': '📝', 'pptx': '📈', 'ppt': '📈',
        'xlsx': '📊', 'xls': '📊', 'csv': '📋', 'json': '🔗', 'xml': '🔗',
        'txt': '📄', 'rtf': '📄', 'md': '📄', 'log': '📄'
    }

    # Document selection with better UI
    st.subheader("📋 Select Documents to Analyze")
    
    doc_ids = list(st.session_state.documents.keys())
    
    # Show available documents with file type info
    st.write("**Available Documents:**")
    available_docs_info = []
    for doc_id in doc_ids:
        doc = st.session_state.documents[doc_id]
        status = "✅ Ready" if not doc['has_error'] else "❌ Error"
        available_docs_info.append({
            'Icon': file_type_icons.get(doc['type'], '📄'),
            'Name': doc['name'],
            'Type': doc['type'].upper(),
            'Words': doc['word_count'],
            'Status': status
        })
    
    df_available = pd.DataFrame(available_docs_info)
    st.dataframe(df_available, use_container_width=True, hide_index=True)
    
    # Multi-select for documents
    selected_docs = st.multiselect(
        "Select documents to compare (minimum 2)",
        options=doc_ids,
        format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}",
        default=doc_ids[:min(3, len(doc_ids))]  # Default to first 3 documents
    )

    if len(selected_docs) >= 2:
        # Show selected documents summary
        st.subheader("🎯 Selected Documents for Analysis")
        selected_docs_info = []
        valid_docs = []
        
        for doc_id in selected_docs:
            doc = st.session_state.documents[doc_id]
            status = "✅ Ready" if not doc['has_error'] else "❌ Error"
            selected_docs_info.append({
                'Icon': file_type_icons.get(doc['type'], '📄'),
                'Name': doc['name'],
                'Type': doc['type'].upper(),
                'Words': doc['word_count'],
                'Status': status
            })
            
            if not doc['has_error']:
                valid_docs.append(doc_id)
        
        df_selected = pd.DataFrame(selected_docs_info)
        st.dataframe(df_selected, use_container_width=True, hide_index=True)

        # Show analysis options
        st.subheader("⚙️ Analysis Configuration")
        
        col1, col2 = st.columns(2)
        with col1:
            analysis_types = st.multiselect(
                "Select analysis types:",
                options=['Temporal', 'Numerical', 'Factual', 'Status', 'Semantic'],
                default=['Temporal', 'Numerical', 'Factual', 'Status', 'Semantic'],
                help="Choose which types of contradictions to detect"
            )
        
        with col2:
            confidence_threshold = st.slider(
                "Minimum confidence threshold:",
                min_value=0.1,
                max_value=1.0,
                value=0.6,
                step=0.1,
                help="Only show contradictions above this confidence level"
            )

        # Analysis button
        if st.button("🔍 Detect Contradictions", type="primary", use_container_width=True):
            if len(valid_docs) < 2:
                st.error("❌ Need at least 2 successfully processed documents for analysis.")
                return
            
            # Prepare documents for analysis
            documents_for_analysis = []
            for doc_id in valid_docs:
                doc = st.session_state.documents[doc_id]
                documents_for_analysis.append({
                    'id': doc_id,
                    'name': doc['name'],
                    'text': doc['text'],
                    'has_error': doc['has_error'],
                    'type': doc['type']
                })

            # Run the enhanced contradiction detection
            with st.spinner("🔍 Analyzing documents for contradictions..."):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                status_text.text("Loading AI models...")
                progress_bar.progress(0.2)
                
                status_text.text("Extracting patterns and entities...")
                progress_bar.progress(0.4)
                
                status_text.text("Detecting contradictions...")
                progress_bar.progress(0.6)
                
                # Use the enhanced detection function
                contradictions = find_contradictions_enhanced(documents_for_analysis)
                
                status_text.text("Processing results...")
                progress_bar.progress(0.8)
                
                # Filter by confidence threshold and selected types
                filtered_contradictions = [
                    c for c in contradictions 
                    if c['confidence'] >= confidence_threshold and 
                    c['type'].title() in analysis_types
                ]
                
                progress_bar.progress(1.0)
                status_text.empty()
                progress_bar.empty()

            # Display results
            st.subheader("📊 Analysis Results")
            
            if filtered_contradictions:
                # Summary metrics
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    st.metric("Total Contradictions", len(filtered_contradictions))
                
                with col2:
                    high_severity = len([c for c in filtered_contradictions if c['severity'] in ['Critical', 'High']])
                    st.metric("High Priority", high_severity)
                
                with col3:
                    avg_confidence = np.mean([c['confidence'] for c in filtered_contradictions])
                    st.metric("Avg Confidence", f"{avg_confidence:.2f}")
                
                with col4:
                    types_found = len(set(c['type'] for c in filtered_contradictions))
                    st.metric("Types Found", types_found)

                # Contradiction breakdown by type
                st.subheader("📈 Contradiction Breakdown")
                contradiction_types = pd.Series([c['type'] for c in filtered_contradictions]).value_counts()
                
                col1, col2 = st.columns([2, 1])
                with col1:
                    fig = px.bar(
                        x=contradiction_types.index, 
                        y=contradiction_types.values,
                        title='Contradictions by Type',
                        labels={'x': 'Type', 'y': 'Count'},
                        color=contradiction_types.values,
                        color_continuous_scale='Reds'
                    )
                    fig.update_xaxes(tickangle=0)
                    st.plotly_chart(fig, use_container_width=True)
                
                with col2:
                    st.write("**Summary by Type:**")
                    for ctype, count in contradiction_types.items():
                        percentage = (count / len(filtered_contradictions)) * 100
                        severity_icon = {
                            'temporal': '⏰', 'numerical': '🔢', 'factual': '📋', 
                            'status': '📊', 'semantic': '🧠'
                        }.get(ctype.lower(), '⚠️')
                        st.write(f"{severity_icon} **{ctype.title()}**: {count} ({percentage:.1f}%)")

                # Detailed contradiction display
                st.subheader("🔍 Detailed Contradictions")
                
                # Sort by severity and confidence
                severity_order = {'Critical': 4, 'High': 3, 'Medium': 2, 'Low': 1}
                sorted_contradictions = sorted(
                    filtered_contradictions,
                    key=lambda x: (severity_order.get(x['severity'], 0), x['confidence']),
                    reverse=True
                )

                for i, contradiction in enumerate(sorted_contradictions):
                    # Severity color coding
                    severity_colors = {
                        'Critical': '🔴', 'High': '🟠', 'Medium': '🟡', 'Low': '🟢'
                    }
                    severity_icon = severity_colors.get(contradiction['severity'], '⚠️')
                    
                    # Type icons
                    type_icons = {
                        'temporal': '⏰', 'numerical': '🔢', 'factual': '📋',
                        'status': '📊', 'semantic': '🧠'
                    }
                    type_icon = type_icons.get(contradiction['type'].lower(), '⚠️')
                    
                    with st.expander(
                        f"{severity_icon} {type_icon} **{contradiction['type'].title()} Contradiction** "
                        f"(Confidence: {contradiction['confidence']:.2f})",
                        expanded=(i < 3)  # Expand first 3 by default
                    ):
                        # Main contradiction info
                        st.write(f"**Description:** {contradiction['description']}")
                        st.write(f"**Severity:** {contradiction['severity']} | **Confidence:** {contradiction['confidence']:.2f}")
                        
                        # Document comparison
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.write(f"**📄 {contradiction['doc1_name']}**")
                            st.info(f"**Value/Statement:** {contradiction['doc1_value']}")
                        
                        with col2:
                            st.write(f"**📄 {contradiction['doc2_name']}**")
                            st.error(f"**Value/Statement:** {contradiction['doc2_value']}")
                        
                        # Additional context
                        if contradiction.get('entities'):
                            st.write(f"**🏷️ Entities Involved:** {', '.join(contradiction['entities'])}")
                        
                        if contradiction.get('context'):
                            with st.expander("📝 Additional Context"):
                                for key, value in contradiction['context'].items():
                                    st.write(f"**{key.replace('_', ' ').title()}:** {value}")
                        
                        # Suggestion
                        if contradiction.get('suggestion'):
                            st.success(f"💡 **Suggestion:** {contradiction['suggestion']}")

                # Export options
                st.subheader("📤 Export Results")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    # JSON export
                    all_contradictions_json = json.dumps({
                        'analysis_timestamp': datetime.now().isoformat(),
                        'documents_analyzed': [doc['name'] for doc in documents_for_analysis],
                        'total_contradictions': len(filtered_contradictions),
                        'confidence_threshold': confidence_threshold,
                        'analysis_types': analysis_types,
                        'contradictions': filtered_contradictions
                    }, indent=2)
                    
                    st.download_button(
                        label="📥 Export as JSON",
                        data=all_contradictions_json,
                        file_name=f"contradiction_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                        mime="application/json"
                    )
                
                with col2:
                    # CSV export
                    csv_data = []
                    for i, c in enumerate(filtered_contradictions):
                        csv_data.append({
                            'ID': i + 1,
                            'Type': c['type'],
                            'Description': c['description'],
                            'Severity': c['severity'],
                            'Confidence': c['confidence'],
                            'Document_1': c['doc1_name'],
                            'Document_1_Value': c['doc1_value'],
                            'Document_2': c['doc2_name'],
                            'Document_2_Value': c['doc2_value'],
                            'Suggestion': c.get('suggestion', ''),
                            'Entities': ', '.join(c.get('entities', []))
                        })
                    
                    df_export = pd.DataFrame(csv_data)
                    csv_string = df_export.to_csv(index=False)
                    
                    st.download_button(
                        label="📊 Export as CSV",
                        data=csv_string,
                        file_name=f"contradiction_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                        mime="text/csv"
                    )
                
                with col3:
                    # Summary report
                    summary_report = f"""# Contradiction Analysis Report

**Analysis Date:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
**Documents Analyzed:** {len(documents_for_analysis)}
**Total Contradictions Found:** {len(filtered_contradictions)}
**Confidence Threshold:** {confidence_threshold}
**Analysis Types:** {', '.join(analysis_types)}

## Summary by Severity:
"""
                    severity_counts = pd.Series([c['severity'] for c in filtered_contradictions]).value_counts()
                    for severity, count in severity_counts.items():
                        summary_report += f"- {severity}: {count}\n"
                    
                    summary_report += f"""
## Summary by Type:
"""
                    for ctype, count in contradiction_types.items():
                        summary_report += f"- {ctype.title()}: {count}\n"
                    
                    summary_report += f"""
## High Priority Contradictions:
"""
                    high_priority = [c for c in filtered_contradictions if c['severity'] in ['Critical', 'High']]
                    for i, c in enumerate(high_priority[:5]):  # Top 5 high priority
                        summary_report += f"{i+1}. {c['type'].title()} - {c['description']}\n"
                    
                    st.download_button(
                        label="📄 Export Summary",
                        data=summary_report,
                        file_name=f"contradiction_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
                        mime="text/markdown"
                    )

            else:
                # No contradictions found
                st.success("🎉 **Excellent! No contradictions found!**")
                st.write("Your documents appear to be consistent with each other.")
                
                # Show what was analyzed
                st.info(f"""
                **Analysis Summary:**
                - Documents analyzed: {len(documents_for_analysis)}
                - Analysis types: {', '.join(analysis_types)}
                - Confidence threshold: {confidence_threshold}
                - Total sentences analyzed: {sum(len(doc['text'].split('.')) for doc in documents_for_analysis)}
                """)
                
                # Suggestions for improvement
                st.write("### 💡 Suggestions:")
                st.write("- Try lowering the confidence threshold to catch more subtle inconsistencies")
                st.write("- Add more documents to expand the analysis scope")
                st.write("- Consider running language enhancement analysis for style consistency")

    else:
        st.warning("⚠️ Please select at least 2 documents to perform contradiction analysis.")
        
        # Show help information
        st.markdown("""
        ### 🎯 How Enhanced Contradiction Detection Works:
        
        **1. Temporal Analysis** ⏰
        - Detects conflicting dates, deadlines, and timelines
        - Example: "Due March 15" vs "Due April 20"
        
        **2. Numerical Analysis** 🔢  
        - Finds conflicting numbers, budgets, quantities
        - Example: "Budget $50,000" vs "Budget $75,000"
        
        **3. Factual Analysis** 📋
        - Identifies contradictory statements about entities
        - Example: "Project is completed" vs "Project is pending"
        
        **4. Status Analysis** 📊
        - Detects conflicting status information
        - Example: "Approved" vs "Rejected"
        
        **5. Semantic Analysis** 🧠
        - Uses AI to find meaning-based contradictions
        - Example: "Cannot be done" vs "Will be completed"
        
        ### 📋 Requirements:
        For best results, install these optional dependencies:
        ```bash
        pip install spacy sentence-transformers
        python -m spacy download en_core_web_sm
        ```
        """)

# =============================================================================
# SUPPORTED FILE TYPES INFO (Keep your existing function)
# =============================================================================

def get_supported_file_types():
    """Return information about supported file types and their availability"""
    return {
        'PDF': {
            'extensions': ['pdf'],
            'available': PDF_AVAILABLE,
            'install_command': 'pip install PyMuPDF',
            'description': 'Portable Document Format'
        },
        'Microsoft Word': {
            'extensions': ['docx', 'doc'],
            'available': DOCX_AVAILABLE,
            'install_command': 'pip install python-docx',
            'description': 'Word documents with text and tables'
        },
        'Microsoft PowerPoint': {
            'extensions': ['pptx', 'ppt'],
            'available': PPTX_AVAILABLE,
            'install_command': 'pip install python-pptx',
            'description': 'PowerPoint presentations with slides and tables'
        },
        'Microsoft Excel (New)': {
            'extensions': ['xlsx'],
            'available': XLSX_AVAILABLE,
            'install_command': 'pip install openpyxl',
            'description': 'Excel workbooks (2007+)'
        },
        'Microsoft Excel (Legacy)': {
            'extensions': ['xls'],
            'available': XLS_AVAILABLE,
            'install_command': 'pip install xlrd',
            'description': 'Excel workbooks (97-2003)'
        },
        'CSV': {
            'extensions': ['csv'],
            'available': True,
            'install_command': 'Built-in',
            'description': 'Comma-separated values'
        },
        'JSON': {
            'extensions': ['json'],
            'available': True,
            'install_command': 'Built-in',
            'description': 'JavaScript Object Notation'
        },
        'RTF': {
            'extensions': ['rtf'],
            'available': RTF_AVAILABLE,
            'install_command': 'pip install striprtf',
            'description': 'Rich Text Format'
        },
        'XML': {
            'extensions': ['xml'],
            'available': XML_AVAILABLE,
            'install_command': 'Built-in',
            'description': 'Extensible Markup Language'
        },
        'Text Files': {
            'extensions': ['txt', 'md', 'log'],
            'available': True,
            'install_command': 'Built-in',
            'description': 'Plain text, Markdown, and log files'
        }
    }

# =============================================================================
# STYLING (Keep your existing styling)
# =============================================================================

st.markdown("""
<style>
    .main {
        padding-top: 2rem;
    }
    .stButton>button {
        background-color: #0078d4;
        color: white;
    }
            
    .header {
        background-color: #f3f2f1;
        padding: 1rem 2rem;
        border-bottom: 1px solid #ccc;
        position: sticky;
        top: 0;
        z-index: 999;
    }        
    body, div, h1, h2, h3, h4, p {
        font-family: 'Segoe UI', sans-serif;
    }
    section[data-testid="stSidebar"] {
        background-color: #ffffff;
        border-right: 1px solid #ccc;
        padding-top: 20px;
    }
            
    section[data-testid="stSidebar"] * {
        color: #262730 !important;
                    
    .stButton>button:hover {
        background-color: #106ebe;
    }
    .success-message {
        padding: 1rem;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 0.25rem;
        color: #155724;
    }
    .error-message {
        padding: 1rem;
        background-color: #f8d7da;
        border: 1px solid #f5c6cb;
        border-radius: 0.25rem;
        color: #721c24;
    }
    .file-support-card {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 0.5rem;
        border: 1px solid #e1e1e1;
    }
    .supported {
        background-color: #d4edda;
        border-color: #c3e6cb;
    }
    .not-supported {
        background-color: #f8d7da;
        border-color: #f5c6cb;
    }
</style>
""", unsafe_allow_html=True)

# =============================================================================
# DOCUMENTATION HTML (Keep your existing documentation)
# =============================================================================

DOCUMENTATION_HTML = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>CoorSight User Documentation</title>
    <style>
        :root {
            --primary-color: #ff4b4b;
            --secondary-color: #ff6b6b;
            --accent-color: #40e0d0;
            --success-color: #00cc88;
            --warning-color: #ffab00;
            --error-color: #ff4b4b;
            --background-color: #fafafa;
            --surface-color: #ffffff;
            --text-primary: #262730;
            --text-secondary: #8e8ea0;
            --border-color: #e6eaf1;
            --shadow: 0 2px 8px rgba(0, 0, 0, 0.1);
            --radius: 8px;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Source Sans Pro', sans-serif;
            background-color: var(--background-color);
            color: var(--text-primary);
            line-height: 1.6;
        }

        .container {
            max-width: 1200px;
            margin: 0 auto;
            padding: 1rem;
        }

        .nav-tabs {
            display: flex;
            background: var(--surface-color);
            border-radius: var(--radius);
            margin-bottom: 2rem;
            box-shadow: var(--shadow);
            overflow-x: auto;
        }

        .nav-tab {
            flex: 1;
            padding: 1rem 1.5rem;
            border: none;
            background: none;
            cursor: pointer;
            font-size: 1rem;
            color: var(--text-secondary);
            transition: all 0.3s ease;
            white-space: nowrap;
            min-width: 120px;
            font-weight: 500;
        }

        .nav-tab.active {
            background: var(--primary-color);
            color: white;
        }

        .nav-tab:hover:not(.active) {
            background: var(--background-color);
            color: var(--primary-color);
        }

        .content-section {
            display: none;
            animation: fadeIn 0.3s ease-in-out;
        }

        .content-section.active {
            display: block;
        }

        @keyframes fadeIn {
            from { opacity: 0; }
            to { opacity: 1; }
        }

        .card {
            background: var(--surface-color);
            border-radius: var(--radius);
            padding: 1.5rem;
            margin-bottom: 1.5rem;
            box-shadow: var(--shadow);
        }

        .card-title {
            font-size: 1.5rem;
            color: var(--primary-color);
            margin-bottom: 1rem;
            font-weight: 600;
        }

        .file-support-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 1rem;
            margin-top: 1rem;
        }

        .file-type-card {
            border: 2px solid var(--border-color);
            border-radius: var(--radius);
            padding: 1rem;
            transition: all 0.3s ease;
        }

        .file-type-card.supported {
            border-color: var(--success-color);
            background: rgba(0, 204, 136, 0.05);
        }

        .file-type-card.not-supported {
            border-color: var(--error-color);
            background: rgba(255, 75, 75, 0.05);
        }

        .file-extensions {
            display: flex;
            gap: 0.5rem;
            margin: 0.5rem 0;
        }

        .extension-tag {
            background: var(--text-secondary);
            color: white;
            padding: 0.25rem 0.5rem;
            border-radius: 4px;
            font-size: 0.8rem;
            font-weight: 600;
        }
    </style>
</head>
<body>
    <div class="container">
        <nav class="nav-tabs">
            <button class="nav-tab active" onclick="showSection('overview')">📋 Overview</button>
            <button class="nav-tab" onclick="showSection('filetypes')">📄 File Types</button>
            <button class="nav-tab" onclick="showSection('quickstart')">⚡ Quick Start</button>
            <button class="nav-tab" onclick="showSection('examples')">💡 Examples</button>
        </nav>

        <!-- Overview Section -->
        <div id="overview" class="content-section active">
            <div class="card">
                <h2 class="card-title">🎯 Enhanced CoorSight Platform</h2>
                <p>Advanced document analysis platform supporting 15+ file formats with AI-powered insights.</p>
                
                <h3 style="margin-top: 2rem;">🌟 What You Can Do</h3>
                <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 1rem; margin-top: 1rem;">
                    <div style="text-align: center; padding: 1rem; background: #f8f9fa; border-radius: 8px;">
                        <div style="font-size: 2rem;">📄</div>
                        <strong>Multi-Format Upload</strong>
                        <p style="font-size: 0.8rem;">PDF, Word, Excel, PowerPoint, CSV, JSON & more</p>
                    </div>
                    <div style="text-align: center; padding: 1rem; background: #f8f9fa; border-radius: 8px;">
                        <div style="font-size: 2rem;">🧠</div>
                        <strong>AI-Powered Analysis</strong>
                        <p style="font-size: 0.8rem;">NLP, sentiment, entities, key phrases</p>
                    </div>
                    <div style="text-align: center; padding: 1rem; background: #f8f9fa; border-radius: 8px;">
                        <div style="font-size: 2rem;">⚖️</div>
                        <strong>Smart Detection</strong>
                        <p style="font-size: 0.8rem;">Contradictions, inconsistencies, quality issues</p>
                    </div>
                    <div style="text-align: center; padding: 1rem; background: #f8f9fa; border-radius: 8px;">
                        <div style="font-size: 2rem;">✨</div>
                        <strong>Quality Enhancement</strong>
                        <p style="font-size: 0.8rem;">Language improvements, style suggestions</p>
                    </div>
                </div>
            </div>
        </div>

        <!-- File Types Section -->
        <div id="filetypes" class="content-section">
            <div class="card">
                <h2 class="card-title">📄 Supported File Types</h2>
                <p>CoorSight supports a wide range of document formats for comprehensive analysis.</p>
                
                <div class="file-support-grid">
                    <div class="file-type-card supported">
                        <h3>📄 PDF Documents</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.pdf</span>
                        </div>
                        <p>Portable Document Format with text extraction</p>
                    </div>

                    <div class="file-type-card supported">
                        <h3>📝 Microsoft Word</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.docx</span>
                            <span class="extension-tag">.doc</span>
                        </div>
                        <p>Word documents with paragraphs and tables</p>
                    </div>

                    <div class="file-type-card supported">
                        <h3>📊 Microsoft Excel</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.xlsx</span>
                            <span class="extension-tag">.xls</span>
                        </div>
                        <p>Excel workbooks with multiple sheets</p>
                    </div>

                    <div class="file-type-card supported">
                        <h3>📈 Microsoft PowerPoint</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.pptx</span>
                            <span class="extension-tag">.ppt</span>
                        </div>
                        <p>PowerPoint presentations with slides</p>
                    </div>

                    <div class="file-type-card supported">
                        <h3>📋 Data Files</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.csv</span>
                            <span class="extension-tag">.json</span>
                            <span class="extension-tag">.xml</span>
                        </div>
                        <p>Structured data formats</p>
                    </div>

                    <div class="file-type-card supported">
                        <h3>📄 Text Files</h3>
                        <div class="file-extensions">
                            <span class="extension-tag">.txt</span>
                            <span class="extension-tag">.rtf</span>
                            <span class="extension-tag">.md</span>
                            <span class="extension-tag">.log</span>
                        </div>
                        <p>Plain text and rich text formats</p>
                    </div>
                </div>
                
                <div style="margin-top: 2rem; padding: 1rem; background: #e8f4fd; border-radius: 8px;">
                    <h4>📦 Installation Requirements</h4>
                    <p>Some file types require additional Python packages:</p>
                    <ul style="margin-top: 0.5rem;">
                        <li><strong>PDF:</strong> <code>pip install PyMuPDF</code></li>
                        <li><strong>Word:</strong> <code>pip install python-docx</code></li>
                        <li><strong>PowerPoint:</strong> <code>pip install python-pptx</code></li>
                        <li><strong>Excel:</strong> <code>pip install openpyxl xlrd</code></li>
                        <li><strong>RTF:</strong> <code>pip install striprtf</code></li>
                    </ul>
                </div>
            </div>
        </div>

        <!-- Quick Start Section -->
        <div id="quickstart" class="content-section">
            <div class="card">
                <h2 class="card-title">⚡ Quick Start Guide</h2>
                
                <div style="display: flex; align-items: center; margin: 2rem 0; padding: 1.5rem; background: #fff; border-left: 4px solid var(--primary-color);">
                    <div style="background: var(--primary-color); color: white; width: 2rem; height: 2rem; border-radius: 50%; display: flex; align-items: center; justify-content: center; margin-right: 1rem;">1</div>
                    <div>
                        <h3>Upload Your Documents</h3>
                        <p>Upload any supported file type including PDF, Word, Excel, PowerPoint, CSV, JSON, and more.</p>
                    </div>
                </div>

                <div style="display: flex; align-items: center; margin: 2rem 0; padding: 1.5rem; background: #fff; border-left: 4px solid var(--primary-color);">
                    <div style="background: var(--primary-color); color: white; width: 2rem; height: 2rem; border-radius: 50%; display: flex; align-items: center; justify-content: center; margin-right: 1rem;">2</div>
                    <div>
                        <h3>Run Analysis</h3>
                        <p>Use the Comprehensive Analysis page to extract insights and analyze content from any file format.</p>
                    </div>
                </div>

                <div style="display: flex; align-items: center; margin: 2rem 0; padding: 1.5rem; background: #fff; border-left: 4px solid var(--primary-color);">
                    <div style="background: var(--primary-color); color: white; width: 2rem; height: 2rem; border-radius: 50%; display: flex; align-items: center; justify-content: center; margin-right: 1rem;">3</div>
                    <div>
                        <h3>View Results</h3>
                        <p>Check the Reports & Insights page for comprehensive analysis results across all document types.</p>
                    </div>
                </div>
            </div>
        </div>

        <!-- Examples Section -->
        <div id="examples" class="content-section">
            <div class="card">
                <h2 class="card-title">💡 File Type Examples</h2>
                
                <h3>📊 Excel Spreadsheet Analysis</h3>
                <p>Upload .xlsx or .xls files to analyze:</p>
                <ul>
                    <li>Financial data and budgets</li>
                    <li>Project timelines and schedules</li>
                    <li>Resource allocation tables</li>
                    <li>Performance metrics and KPIs</li>
                </ul>
                
                <h3 style="margin-top: 2rem;">📈 PowerPoint Presentation Analysis</h3>
                <p>Upload .pptx or .ppt files to extract:</p>
                <ul>
                    <li>Slide content and titles</li>
                    <li>Key messaging and themes</li>
                    <li>Data from embedded tables</li>
                    <li>Project updates and status reports</li>
                </ul>

                <h3 style="margin-top: 2rem;">📋 CSV Data Processing</h3>
                <p>Upload .csv files to analyze:</p>
                <ul>
                    <li>Survey responses and feedback</li>
                    <li>Sales data and customer information</li>
                    <li>Log files and system data</li>
                    <li>Research data and measurements</li>
                </ul>

                <h3 style="margin-top: 2rem;">🔗 JSON Data Analysis</h3>
                <p>Upload .json files to process:</p>
                <ul>
                    <li>API responses and configuration files</li>
                    <li>Structured data exports</li>
                    <li>Application settings and metadata</li>
                    <li>Research datasets and surveys</li>
                </ul>
            </div>
        </div>
    </div>

    <script>
        function showSection(sectionName) {
            // Update nav tabs
            document.querySelectorAll('.nav-tab').forEach(tab => {
                tab.classList.remove('active');
            });
            event.target.classList.add('active');

            // Update content sections
            document.querySelectorAll('.content-section').forEach(section => {
                section.classList.remove('active');
            });
            document.getElementById(sectionName).classList.add('active');
        }
    </script>
</body>
</html>
'''

# =============================================================================
# HEADER (Keep your existing header)
# =============================================================================

st.markdown('<div class="header">', unsafe_allow_html=True)
col1, col2, col3 = st.columns([1, 3, 1])

with col2:
    logo_col, text_col = st.columns([1, 5])
    with logo_col:
        st.image("C:/Users/AbigailLarch/Pictures/CoorSightLogo.png", width=80)
    with text_col:
        st.markdown("### CoorSight Document Analysis Platform")
        st.markdown("*Enhanced Multi-Format Document Intelligence with AI Contradiction Detection*")

with col3:
    st.markdown("**Paul Sherman**", unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# =============================================================================
# SIDEBAR NAVIGATION (Keep your existing sidebar)
# =============================================================================

st.sidebar.title("Navigation")

# Enhanced sidebar with AI status indicators
def install_requirements_helper():
    """Helper function to guide users through installing requirements"""
    st.sidebar.markdown("### 🔧 AI Enhancement Status")
    
    if SPACY_AVAILABLE:
        st.sidebar.success("✅ spaCy NLP Ready")
    else:
        st.sidebar.warning("⚠️ spaCy Missing")
        st.sidebar.code("pip install spacy")
        st.sidebar.code("python -m spacy download en_core_web_sm")
    
    if SENTENCE_TRANSFORMER_AVAILABLE:
        st.sidebar.success("✅ Semantic Analysis Ready")
    else:
        st.sidebar.warning("⚠️ Sentence Transformers Missing")
        st.sidebar.code("pip install sentence-transformers")
    
    try:
        from enhanced_contradiction_detection import EnhancedContradictionDetector
        st.sidebar.success("✅ Enhanced Detection Module Ready")
    except ImportError:
        st.sidebar.error("❌ Enhanced Detection Module Missing")
        st.sidebar.write("Save the enhanced_contradiction_detection.py file in your project directory")

# Show the helper
install_requirements_helper()

pages = [
    "📤 Upload Documents",
    "📚 Document Library", 
    "🔍 Comprehensive Analysis",
    "⚖️ Contradiction Detection",  # Updated with enhanced symbol
    "✨ Language Enhancement",
    "🔀 Document Comparison",
    "📈 Reports & Insights",
    "📖 User Documentation",
    "💚 System Health"
]
selected_page = st.sidebar.radio("Go to", pages, label_visibility="collapsed")
st.session_state.current_page = selected_page

# =============================================================================
# ANALYSIS FUNCTIONS (Keep your existing analysis functions)
# =============================================================================

def analyze_nlp(text: str) -> Dict[str, Any]:
    words = text.split()
    sentences = text.split('.')
    entities = [
        {'text': match.group(), 'type': 'PERSON', 'confidence': 0.7}
        for match in re.finditer(r'\b[A-Z][a-z]+ [A-Z][a-z]+\b', text)
    ]
    word_freq = {}
    for word in words:
        if len(word) > 4:
            word_freq[word.lower()] = word_freq.get(word.lower(), 0) + 1
    key_phrases = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
    pos_words = ['good', 'great', 'excellent', 'positive', 'success']
    neg_words = ['bad', 'poor', 'negative', 'failure', 'problem']
    pos_count = sum(1 for word in words if word.lower() in pos_words)
    neg_count = sum(1 for word in words if word.lower() in neg_words)
    sentiment = "Positive" if pos_count > neg_count else "Negative" if neg_count > pos_count else "Neutral"
    return {
        'word_count': len(words),
        'sentence_count': len(sentences),
        'entities': entities[:5],
        'key_phrases': [{'text': word, 'score': count} for word, count in key_phrases],
        'sentiment': sentiment
    }

import re
from collections import Counter
from typing import Dict, Any, List, Optional
from enum import Enum
import logging

class DocumentType(Enum):
    """Supported document types with different analysis profiles"""
    TECHNICAL = "technical"
    BUSINESS = "business" 
    MARKETING = "marketing"
    LEGAL = "legal"
    EMAIL = "email"
    REPORT = "report"
    POLICY = "policy"
    GENERAL = "general"

class AnalysisProfile:
    """Configuration profiles for different document types"""
    
    PROFILES = {
        DocumentType.TECHNICAL: {
            'max_sentence_length': 30,  # Technical docs can be longer
            'max_avg_chars_per_word': 7,  # Technical terms are longer
            'passive_voice_threshold': 0.4,  # More acceptable in technical writing
            'paragraph_word_threshold': 200,
            'check_jargon': True,
            'severity_multiplier': 0.8  # Less strict
        },
        DocumentType.BUSINESS: {
            'max_sentence_length': 25,
            'max_avg_chars_per_word': 6,
            'passive_voice_threshold': 0.25,
            'paragraph_word_threshold': 150,
            'check_jargon': True,
            'severity_multiplier': 1.0
        },
        DocumentType.MARKETING: {
            'max_sentence_length': 20,  # Punchier sentences
            'max_avg_chars_per_word': 5.5,
            'passive_voice_threshold': 0.15,  # Active voice preferred
            'paragraph_word_threshold': 100,
            'check_jargon': False,  # Marketing can use buzzwords
            'severity_multiplier': 1.2  # More strict for readability
        },
        DocumentType.LEGAL: {
            'max_sentence_length': 35,  # Legal language can be complex
            'max_avg_chars_per_word': 8,
            'passive_voice_threshold': 0.5,  # Common in legal writing
            'paragraph_word_threshold': 250,
            'check_jargon': False,  # Legal jargon is necessary
            'severity_multiplier': 0.6  # Less strict on style
        },
        DocumentType.EMAIL: {
            'max_sentence_length': 18,
            'max_avg_chars_per_word': 5.5,
            'passive_voice_threshold': 0.2,
            'paragraph_word_threshold': 80,
            'check_jargon': True,
            'severity_multiplier': 1.1
        },
        DocumentType.GENERAL: {
            'max_sentence_length': 25,
            'max_avg_chars_per_word': 6,
            'passive_voice_threshold': 0.3,
            'paragraph_word_threshold': 150,
            'check_jargon': True,
            'severity_multiplier': 1.0
        }
    }

def analyze_document_language(
    text: str, 
    document_type: DocumentType = DocumentType.GENERAL,
    document_id: Optional[str] = None,
    source_system: Optional[str] = None
) -> Dict[str, Any]:
    """
    Enhanced language analysis function for document management systems.
    
    Args:
        text: The text content to analyze
        document_type: Type of document for tailored analysis
        document_id: Optional document identifier for logging
        source_system: Optional source system identifier
    
    Returns:
        Dictionary with suggestions, metrics, and metadata
    """
    
    # Initialize logging context
    analysis_context = {
        'document_id': document_id,
        'source_system': source_system,
        'document_type': document_type.value
    }
    
    # Input validation
    if not text or not text.strip():
        return {
            'suggestions': [{
                'type': 'Error',
                'issue': 'Empty or null text content',
                'suggestion': 'Document contains no analyzable text content',
                'severity': 'High',
                'category': 'data_quality'
            }],
            'metrics': _get_empty_metrics(),
            'metadata': analysis_context,
            'analysis_score': 0
        }
    
    # Get analysis profile for document type
    profile = AnalysisProfile.PROFILES.get(document_type, AnalysisProfile.PROFILES[DocumentType.GENERAL])
    
    suggestions = []
    text = text.strip()
    
    # Basic metrics calculation
    metrics = _calculate_basic_metrics(text)
    
    # Perform analysis checks based on document type
    suggestions.extend(_check_sentence_structure(text, metrics, profile))
    suggestions.extend(_check_passive_voice(text, metrics, profile))
    suggestions.extend(_check_word_usage(text, metrics, profile))
    suggestions.extend(_check_paragraph_structure(text, metrics, profile))
    suggestions.extend(_check_readability(text, metrics, profile))
    suggestions.extend(_check_grammar_issues(text, profile))
    
    # Document type specific checks
    if document_type == DocumentType.TECHNICAL:
        suggestions.extend(_check_technical_writing(text, metrics))
    elif document_type == DocumentType.MARKETING:
        suggestions.extend(_check_marketing_copy(text, metrics))
    elif document_type == DocumentType.EMAIL:
        suggestions.extend(_check_email_format(text, metrics))
    
    # Calculate overall analysis score
    analysis_score = _calculate_analysis_score(metrics, suggestions, profile)
    
    # Apply severity adjustments based on document type
    suggestions = _adjust_severity_by_profile(suggestions, profile)
    
    return {
        'suggestions': suggestions,
        'metrics': metrics,
        'metadata': {
            **analysis_context,
            'analysis_profile': document_type.value,
            'total_issues': len(suggestions),
            'high_severity_issues': len([s for s in suggestions if s['severity'] == 'High'])
        },
        'analysis_score': analysis_score
    }

def _calculate_basic_metrics(text: str) -> Dict[str, Any]:
    """Calculate basic text metrics"""
    words = text.split()
    word_count = len(words)
    char_count = len(text)
    
    # Improved sentence splitting
    sentences = re.split(r'[.!?]+(?:\s+|$)', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    sentence_count = len(sentences)
    
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    paragraph_count = len(paragraphs)
    
    avg_words_per_sentence = word_count / sentence_count if sentence_count > 0 else 0
    avg_chars_per_word = char_count / word_count if word_count > 0 else 0
    
    # Calculate passive voice ratio
    passive_patterns = [
        r'\b(was|were|is|are|been|be)\s+\w*ed\b',
        r'\b(was|were|is|are|been|be)\s+\w+en\b'
    ]
    passive_count = sum(len(re.findall(pattern, text, re.IGNORECASE)) for pattern in passive_patterns)
    passive_ratio = passive_count / sentence_count if sentence_count > 0 else 0
    
    return {
        'word_count': word_count,
        'sentence_count': sentence_count,
        'paragraph_count': paragraph_count,
        'char_count': char_count,
        'avg_words_per_sentence': round(avg_words_per_sentence, 1),
        'avg_chars_per_word': round(avg_chars_per_word, 1),
        'passive_voice_ratio': round(passive_ratio, 2),
        'passive_count': passive_count,
        'sentences': sentences,
        'words': words
    }

def _check_sentence_structure(text: str, metrics: Dict, profile: Dict) -> List[Dict]:
    """Check sentence length and structure issues"""
    suggestions = []
    sentences = metrics['sentences']
    max_length = profile['max_sentence_length']
    
    long_sentences = []
    very_short_sentences = []
    
    for i, sentence in enumerate(sentences):
        sentence_words = len(sentence.split())
        if sentence_words > max_length:
            long_sentences.append(i + 1)
        elif sentence_words < 4 and sentence_words > 0:
            very_short_sentences.append(i + 1)
    
    if long_sentences:
        severity = 'High' if len(long_sentences) > 3 else 'Medium'
        suggestions.append({
            'type': 'Clarity',
            'issue': f'Long sentences detected in positions: {", ".join(map(str, long_sentences[:5]))}',
            'suggestion': f'Break sentences longer than {max_length} words for better readability',
            'severity': severity,
            'category': 'sentence_structure',
            'affected_sentences': long_sentences
        })
    
    if len(very_short_sentences) > metrics['sentence_count'] * 0.3:
        suggestions.append({
            'type': 'Style',
            'issue': f'Many very short sentences ({len(very_short_sentences)} found)',
            'suggestion': 'Consider combining short sentences for better flow',
            'severity': 'Low',
            'category': 'sentence_structure'
        })
    
    return suggestions

def _check_passive_voice(text: str, metrics: Dict, profile: Dict) -> List[Dict]:
    """Check passive voice usage"""
    suggestions = []
    passive_ratio = metrics['passive_voice_ratio']
    threshold = profile['passive_voice_threshold']
    
    if passive_ratio > threshold:
        severity = 'High' if passive_ratio > threshold * 1.5 else 'Medium'
        suggestions.append({
            'type': 'Style',
            'issue': f'High passive voice usage ({metrics["passive_count"]} instances, {passive_ratio:.1%})',
            'suggestion': 'Use active voice for clearer, more direct communication',
            'severity': severity,
            'category': 'voice'
        })
    
    return suggestions

def _check_word_usage(text: str, metrics: Dict, profile: Dict) -> List[Dict]:
    """Check word usage patterns and repetition"""
    suggestions = []
    words = metrics['words']
    
    # Weak intensifiers check
    weak_intensifiers = ['very', 'really', 'quite', 'rather', 'pretty', 'fairly', 
                        'somewhat', 'extremely', 'absolutely', 'totally', 'completely']
    intensifier_issues = []
    
    for intensifier in weak_intensifiers:
        pattern = r'\b' + re.escape(intensifier) + r'\b'
        matches = re.findall(pattern, text, re.IGNORECASE)
        if matches:
            intensifier_issues.append(f'"{intensifier}" ({len(matches)}x)')
    
    if intensifier_issues:
        severity = 'Medium' if len(intensifier_issues) > 3 else 'Low'
        suggestions.append({
            'type': 'Style',
            'issue': f'Weak intensifiers found: {", ".join(intensifier_issues[:3])}',
            'suggestion': 'Use stronger, more specific adjectives instead of weak intensifiers',
            'severity': severity,
            'category': 'word_choice'
        })
    
    # Word repetition analysis
    content_words = [
        word.lower().strip('.,!?;:"()[]') for word in words 
        if len(word) > 3 and word.lower() not in [
            'that', 'this', 'with', 'from', 'they', 'have', 'will', 
            'been', 'were', 'said', 'than', 'them', 'what', 'when', 
            'where', 'which', 'while', 'would', 'could', 'should'
        ]
    ]
    
    word_freq = Counter(content_words)
    threshold = max(2, metrics['sentence_count'] // 4)
    repeated_words = [(word, count) for word, count in word_freq.items() if count > threshold]
    
    if repeated_words:
        top_repeated = repeated_words[:3]
        word_list = [f'"{word}" ({count}x)' for word, count in top_repeated]
        suggestions.append({
            'type': 'Style',
            'issue': f'Repeated words: {", ".join(word_list)}',
            'suggestion': 'Use synonyms or rephrase to avoid repetition',
            'severity': 'Low',
            'category': 'repetition'
        })
    
    return suggestions

def _check_paragraph_structure(text: str, metrics: Dict, profile: Dict) -> List[Dict]:
    """Check paragraph structure and organization"""
    suggestions = []
    threshold = profile['paragraph_word_threshold']
    
    if metrics['paragraph_count'] == 1 and metrics['word_count'] > threshold:
        suggestions.append({
            'type': 'Structure',
            'issue': f'Single large paragraph ({metrics["word_count"]} words)',
            'suggestion': f'Break text into multiple paragraphs (recommended: <{threshold} words each)',
            'severity': 'Medium',
            'category': 'structure'
        })
    
    return suggestions

def _check_readability(text: str, metrics: Dict, profile: Dict) -> List[Dict]:
    """Check readability metrics"""
    suggestions = []
    
    # Sentence length check
    if metrics['avg_words_per_sentence'] > profile['max_sentence_length']:
        suggestions.append({
            'type': 'Clarity',
            'issue': f'High average sentence length ({metrics["avg_words_per_sentence"]} words)',
            'suggestion': f'Aim for {profile["max_sentence_length"]} words or fewer per sentence',
            'severity': 'Medium',
            'category': 'readability'
        })
    
    # Word complexity check
    if metrics['avg_chars_per_word'] > profile['max_avg_chars_per_word']:
        suggestions.append({
            'type': 'Clarity',
            'issue': f'Complex vocabulary (avg {metrics["avg_chars_per_word"]} chars/word)',
            'suggestion': 'Consider using simpler words where appropriate',
            'severity': 'Low',
            'category': 'readability'
        })
    
    return suggestions

def _check_grammar_issues(text: str, profile: Dict) -> List[Dict]:
    """Check for common grammar and usage errors"""
    suggestions = []
    
    # Common errors
    error_patterns = {
        r'\birregardless\b': 'Use "regardless" instead of "irregardless"',
        r'\balot\b': 'Use "a lot" (two words) instead of "alot"',
        r'\b(should|could|would) of\b': 'Use "should have" instead of "should of"',
        r'\bits\s+it\'s\b': 'Check usage: "its" (possessive) vs "it\'s" (it is)',
        r'\byour\s+you\'re\b': 'Check usage: "your" (possessive) vs "you\'re" (you are)'
    }
    
    found_errors = []
    for pattern, correction in error_patterns.items():
        if re.search(pattern, text, re.IGNORECASE):
            found_errors.append(correction)
    
    if found_errors:
        suggestions.append({
            'type': 'Grammar',
            'issue': 'Common grammar/usage errors detected',
            'suggestion': '; '.join(found_errors[:3]),
            'severity': 'High',
            'category': 'grammar'
        })
    
    return suggestions

def _check_technical_writing(text: str, metrics: Dict) -> List[Dict]:
    """Specific checks for technical documents"""
    suggestions = []
    
    # Check for undefined acronyms (simplified check)
    acronyms = re.findall(r'\b[A-Z]{2,}\b', text)
    if len(set(acronyms)) > 3:
        suggestions.append({
            'type': 'Technical',
            'issue': f'Multiple acronyms detected ({len(set(acronyms))} unique)',
            'suggestion': 'Ensure all acronyms are defined on first use',
            'severity': 'Low',
            'category': 'technical_style'
        })
    
    return suggestions

def _check_marketing_copy(text: str, metrics: Dict) -> List[Dict]:
    """Specific checks for marketing content"""
    suggestions = []
    
    # Check for call-to-action presence in longer marketing copy
    if metrics['word_count'] > 100:
        cta_patterns = [r'\b(click|buy|order|purchase|subscribe|download|register|sign up)\b']
        has_cta = any(re.search(pattern, text, re.IGNORECASE) for pattern in cta_patterns)
        
        if not has_cta:
            suggestions.append({
                'type': 'Marketing',
                'issue': 'No clear call-to-action detected',
                'suggestion': 'Consider adding a clear call-to-action',
                'severity': 'Medium',
                'category': 'marketing_effectiveness'
            })
    
    return suggestions

def _check_email_format(text: str, metrics: Dict) -> List[Dict]:
    """Specific checks for email content"""
    suggestions = []
    
    # Check for appropriate greeting/closing
    has_greeting = bool(re.search(r'^(hi|hello|dear|greetings)', text.strip(), re.IGNORECASE))
    has_closing = bool(re.search(r'(regards|sincerely|thanks|best)', text.strip(), re.IGNORECASE))
    
    if not has_greeting and metrics['word_count'] > 20:
        suggestions.append({
            'type': 'Email',
            'issue': 'No greeting detected',
            'suggestion': 'Consider adding a greeting for professional communication',
            'severity': 'Low',
            'category': 'email_format'
        })
    
    if not has_closing and metrics['word_count'] > 50:
        suggestions.append({
            'type': 'Email',
            'issue': 'No closing detected',
            'suggestion': 'Consider adding a professional closing',
            'severity': 'Low',
            'category': 'email_format'
        })
    
    return suggestions

def _calculate_analysis_score(metrics: Dict, suggestions: List[Dict], profile: Dict) -> int:
    """Calculate overall analysis score (0-100)"""
    base_score = 100
    
    # Penalize based on suggestion severity
    for suggestion in suggestions:
        if suggestion['severity'] == 'High':
            base_score -= 15
        elif suggestion['severity'] == 'Medium':
            base_score -= 8
        else:  # Low
            base_score -= 3
    
    # Apply profile-specific adjustments
    multiplier = profile.get('severity_multiplier', 1.0)
    
    return max(0, min(100, int(base_score * multiplier)))

def _adjust_severity_by_profile(suggestions: List[Dict], profile: Dict) -> List[Dict]:
    """Adjust suggestion severities based on document type profile"""
    multiplier = profile.get('severity_multiplier', 1.0)
    
    if multiplier == 1.0:
        return suggestions
    
    # Adjust severities for less/more strict profiles
    for suggestion in suggestions:
        if multiplier < 1.0:  # Less strict
            if suggestion['severity'] == 'High':
                suggestion['severity'] = 'Medium'
            elif suggestion['severity'] == 'Medium' and multiplier < 0.7:
                suggestion['severity'] = 'Low'
        elif multiplier > 1.0:  # More strict
            if suggestion['severity'] == 'Low':
                suggestion['severity'] = 'Medium'
            elif suggestion['severity'] == 'Medium' and multiplier > 1.3:
                suggestion['severity'] = 'High'
    
    return suggestions

def _get_empty_metrics() -> Dict[str, Any]:
    """Return empty metrics structure"""
    return {
        'word_count': 0,
        'sentence_count': 0,
        'paragraph_count': 0,
        'char_count': 0,
        'avg_words_per_sentence': 0,
        'avg_chars_per_word': 0,
        'passive_voice_ratio': 0,
        'passive_count': 0
    }

# Example usage for document management system integration
def batch_analyze_documents(documents: List[Dict]) -> Dict[str, Any]:
    """
    Analyze multiple documents from a document management system
    
    Args:
        documents: List of dicts with keys: 'text', 'type', 'id', 'source'
    
    Returns:
        Dictionary with individual results and aggregate statistics
    """
    results = {}
    aggregate_stats = {
        'total_documents': len(documents),
        'avg_score': 0,
        'common_issues': Counter(),
        'by_document_type': {}
    }
    
    scores = []
    
    for doc in documents:
        doc_type = DocumentType(doc.get('type', 'general'))
        result = analyze_document_language(
            text=doc['text'],
            document_type=doc_type,
            document_id=doc.get('id'),
            source_system=doc.get('source')
        )
        
        results[doc['id']] = result
        scores.append(result['analysis_score'])
        
        # Aggregate statistics
        for suggestion in result['suggestions']:
            aggregate_stats['common_issues'][suggestion['issue']] += 1
        
        doc_type_key = doc_type.value
        if doc_type_key not in aggregate_stats['by_document_type']:
            aggregate_stats['by_document_type'][doc_type_key] = {
                'count': 0, 'avg_score': 0, 'scores': []
            }
        
        aggregate_stats['by_document_type'][doc_type_key]['count'] += 1
        aggregate_stats['by_document_type'][doc_type_key]['scores'].append(result['analysis_score'])
    
    # Calculate averages
    if scores:
        aggregate_stats['avg_score'] = sum(scores) / len(scores)
        
        for doc_type_stats in aggregate_stats['by_document_type'].values():
            if doc_type_stats['scores']:
                doc_type_stats['avg_score'] = sum(doc_type_stats['scores']) / len(doc_type_stats['scores'])
    
    return {
        'individual_results': results,
        'aggregate_statistics': aggregate_stats
    }

# =============================================================================
# PAGE IMPLEMENTATIONS (Keep all your existing page functions)
# =============================================================================

def page_documentation():
    """Enhanced Documentation Page"""
    st.header("📖 User Documentation")
    st.write("Interactive documentation for the Enhanced CoorSight Document Analysis Platform")
    
    # Embed the HTML documentation
    components.html(DOCUMENTATION_HTML, height=800, scrolling=True)

def page_upload():
    st.header("📤 Upload Documents")
    st.write("Upload your documents for analysis. **Enhanced support for 15+ file formats!**")

    # Show supported file types
    with st.expander("📋 Supported File Types & Status", expanded=False):
        file_types = get_supported_file_types()
        
        for file_type, info in file_types.items():
            status_emoji = "✅" if info['available'] else "❌"
            extensions = ", ".join([f".{ext}" for ext in info['extensions']])
            
            if info['available']:
                st.markdown(f"""
                <div class="file-support-card supported">
                    <strong>{status_emoji} {file_type}</strong><br>
                    <small>Extensions: {extensions}</small><br>
                    <small>{info['description']}</small>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="file-support-card not-supported">
                    <strong>{status_emoji} {file_type}</strong><br>
                    <small>Extensions: {extensions}</small><br>
                    <small>{info['description']}</small><br>
                    <small><code>{info['install_command']}</code></small>
                </div>
                """, unsafe_allow_html=True)

    # Get all supported extensions
    file_types = get_supported_file_types()
    all_extensions = []
    for info in file_types.values():
        all_extensions.extend(info['extensions'])

    uploaded_files = st.file_uploader(
        "Choose files",
        type=all_extensions,
        accept_multiple_files=True,
        help=f"Supported formats: {', '.join(all_extensions)}"
    )

    if uploaded_files:
        if st.button("Process Files", type="primary"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            processed_files = []
            failed_files = []

            for i, uploaded_file in enumerate(uploaded_files):
                status_text.text(f"Processing {uploaded_file.name}...")
                document = process_file(uploaded_file)
                
                if document['has_error']:
                    failed_files.append(document)
                else:
                    processed_files.append(document)
                    st.session_state.documents[document['id']] = document
                
                progress_bar.progress((i + 1) / len(uploaded_files))

            status_text.empty()
            
            # Show results
            if processed_files:
                st.success(f"✅ Successfully processed {len(processed_files)} files!")
                
                st.subheader("📁 Successfully Processed Documents")
                for doc in processed_files:
                    with st.expander(f"📄 {doc['name']} ({doc['type'].upper()})"):
                        col1, col2, col3, col4 = st.columns(4)
                        col1.metric("Word Count", doc['word_count'])
                        col2.metric("Size", f"{doc['size'] / 1024:.1f} KB")
                        col3.metric("Type", doc['type'].upper())
                        col4.metric("Status", "✅ Ready")

                        # Show text preview
                        if doc['text'] and len(doc['text']) > 100:
                            st.text_area("Text Preview", doc['text'][:300] + "...", height=100, key=f"preview_{doc['id']}")

                        if st.button(f"Analyze Now", key=f"analyze_{doc['id']}"):
                            st.session_state.current_page = "🔍 Comprehensive Analysis"
                            st.rerun()
            
            if failed_files:
                st.error(f"❌ Failed to process {len(failed_files)} files")
                with st.expander("View Processing Errors"):
                    for doc in failed_files:
                        st.write(f"**{doc['name']}:** {doc['text']}")

def page_library():
    st.header("📚 Document Library")

    if not st.session_state.documents:
        st.info("No documents uploaded yet. Go to Upload Documents to add files.")
        return

    # Enhanced document display with file type icons
    file_type_icons = {
        'pdf': '📄', 'docx': '📝', 'doc': '📝', 'pptx': '📈', 'ppt': '📈',
        'xlsx': '📊', 'xls': '📊', 'csv': '📋', 'json': '🔗', 'xml': '🔗',
        'txt': '📄', 'rtf': '📄', 'md': '📄', 'log': '📄'
    }

    docs_data = [{
        'Icon': file_type_icons.get(doc['type'], '📄'),
        'Name': doc['name'],
        'Type': doc['type'].upper(),
        'Words': doc['word_count'],
        'Size (KB)': f"{doc['size'] / 1024:.1f}",
        'Uploaded': doc['uploaded_at'].strftime('%Y-%m-%d %H:%M'),
        'ID': doc['id']
    } for doc in st.session_state.documents.values()]

    df = pd.DataFrame(docs_data)
    st.dataframe(df[['Icon', 'Name', 'Type', 'Words', 'Size (KB)', 'Uploaded']], use_container_width=True, hide_index=True)

    # File type statistics
    st.subheader("📊 Library Statistics")
    col1, col2, col3, col4 = st.columns(4)
    
    file_types = [doc['type'] for doc in st.session_state.documents.values()]
    type_counts = pd.Series(file_types).value_counts()
    
    col1.metric("Total Documents", len(st.session_state.documents))
    col2.metric("File Types", len(type_counts))
    col3.metric("Total Words", sum(doc['word_count'] for doc in st.session_state.documents.values()))
    col4.metric("Total Size", f"{sum(doc['size'] for doc in st.session_state.documents.values()) / 1024:.1f} KB")

    # File type breakdown
    if len(type_counts) > 1:
        st.subheader("📈 File Type Distribution")
        fig = px.bar(x=type_counts.index, y=type_counts.values, 
             title='Document Type Distribution',
             labels={'x': 'File Type', 'y': 'Count'})
        fig.update_xaxes(tickangle=0)
        st.plotly_chart(fig, use_container_width=True)

    st.subheader("Document Actions")
    selected_doc = st.selectbox(
        "Select a document",
        options=list(st.session_state.documents.keys()),
        format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}"
    )

    if selected_doc:
        col1, col2, col3 = st.columns(3)
        if col1.button("View Details"):
            doc = st.session_state.documents[selected_doc]
            st.write("**Full Document Text:**")
            st.text_area("", doc['text'], height=300, key="doc_view")

        if col2.button("Analyze"):
            st.session_state.current_page = "🔍 Comprehensive Analysis"
            st.rerun()

        if col3.button("Delete"):
            del st.session_state.documents[selected_doc]
            st.rerun()

def page_comprehensive_analysis():
    st.header("🔍 Comprehensive Analysis")
    st.write("Analyze documents from any supported format with AI-powered insights")

    if not st.session_state.documents:
        st.info("No documents to analyze. Please upload documents first.")
        return

    file_type_icons = {
        'pdf': '📄', 'docx': '📝', 'doc': '📝', 'pptx': '📈', 'ppt': '📈',
        'xlsx': '📊', 'xls': '📊', 'csv': '📋', 'json': '🔗', 'xml': '🔗',
        'txt': '📄', 'rtf': '📄', 'md': '📄', 'log': '📄'
    }

    selected_doc_id = st.selectbox(
        "Select a document to analyze",
        options=list(st.session_state.documents.keys()),
        format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}"
    )

    if selected_doc_id:
        doc = st.session_state.documents[selected_doc_id]
        
        # Show document info
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("File Type", doc['type'].upper())
        col2.metric("Word Count", doc['word_count'])
        col3.metric("Size", f"{doc['size'] / 1024:.1f} KB")
        col4.metric("Uploaded", doc['uploaded_at'].strftime('%m/%d'))

    if st.button("Run Analysis", type="primary"):
        doc = st.session_state.documents[selected_doc_id]

        with st.spinner("Analyzing document..."):
            if doc['text'] and not doc['has_error']:
                nlp_results = analyze_nlp(doc['text'])
                language_results = analyze_document_language(doc['text'])

                analysis_id = str(uuid.uuid4())
                st.session_state.analyses[analysis_id] = {
                    'doc_id': selected_doc_id,
                    'doc_name': doc['name'],
                    'doc_type': doc['type'],
                    'nlp': nlp_results,
                    'language': language_results,
                    'timestamp': datetime.now()
                }

                st.success("Analysis complete!")

                st.subheader("📊 NLP Analysis Results")
                col1, col2, col3 = st.columns(3)
                col1.metric("Word Count", nlp_results['word_count'])
                col2.metric("Sentences", nlp_results['sentence_count'])
                col3.metric("Sentiment", nlp_results['sentiment'])

                st.subheader("🔑 Key Phrases")
                phrases_df = pd.DataFrame(nlp_results['key_phrases'])
                if not phrases_df.empty:
                    fig = px.bar(phrases_df, x='text', y='score', title='Key Phrases')
                    fig.update_xaxes(tickangle=0)
                    st.plotly_chart(fig, use_container_width=True)

                st.subheader("✨ Language Quality")
                col1, col2 = st.columns(2)
                col1.metric(f"{language_results}/100")
                col2.metric("Avg Words/Sentence", f"{language_results['metrics']['avg_words_per_sentence']:.1f}")

                if language_results['suggestions']:
                    st.write("**Improvement Suggestions:**")
                    for sugg in language_results['suggestions']:
                        st.warning(f"{sugg['type']}: {sugg['suggestion']}")
            else:
                st.error("Cannot analyze this document due to processing errors.")

def page_language_enhancement():
    st.header("✨ Language Enhancement")
    st.write("Improve writing quality and clarity across all document types")

    if not st.session_state.documents:
        st.info("No documents to analyze. Please upload documents first.")
        return

    file_type_icons = {
        'pdf': '📄', 'docx': '📝', 'doc': '📝', 'pptx': '📈', 'ppt': '📈',
        'xlsx': '📊', 'xls': '📊', 'csv': '📋', 'json': '🔗', 'xml': '🔗',
        'txt': '📄', 'rtf': '📄', 'md': '📄', 'log': '📄'
    }

    selected_doc_id = st.selectbox(
        "Select a document for language analysis",
        options=list(st.session_state.documents.keys()),
        format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}"
    )

    if st.button("Analyze Language", type="primary"):
        doc = st.session_state.documents[selected_doc_id]

        if doc['has_error']:
            st.error("Cannot analyze this document due to processing errors.")
            return

        with st.spinner("Analyzing language..."):
            results = analyze_document_language(doc['text'])

        col1, col2, col3 = st.columns(3)
        col2.metric("Avg Words/Sentence", f"{results['metrics']['avg_words_per_sentence']:.1f}")
        col3.metric("Total Sentences", results['metrics']['sentence_count'])

        if results['suggestions']:
            st.subheader("📝 Improvement Suggestions")
            for sugg in results['suggestions']:
                with st.expander(f"{sugg['type']}: {sugg['issue']}"):
                    st.write(sugg['suggestion'])
                    st.caption(f"Severity: {sugg['severity']}")
        else:
            st.success("Great job! No major language issues found.")

def page_comparison():
    st.header("🔀 Document Comparison")
    st.write("Compare documents across different formats and find similarities and differences")

    if len(st.session_state.documents) < 2:
        st.info("Need at least 2 documents to compare.")
        return

    file_type_icons = {
        'pdf': '📄', 'docx': '📝', 'doc': '📝', 'pptx': '📈', 'ppt': '📈',
        'xlsx': '📊', 'xls': '📊', 'csv': '📋', 'json': '🔗', 'xml': '🔗',
        'txt': '📄', 'rtf': '📄', 'md': '📄', 'log': '📄'
    }

    col1, col2 = st.columns(2)

    with col1:
        doc1_id = st.selectbox(
            "Document A",
            options=list(st.session_state.documents.keys()),
            format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}",
            key="doc_a"
        )

    with col2:
        doc2_id = st.selectbox(
            "Document B",
            options=list(st.session_state.documents.keys()),
            format_func=lambda x: f"{file_type_icons.get(st.session_state.documents[x]['type'], '📄')} {st.session_state.documents[x]['name']}",
            key="doc_b"
        )

    if doc1_id and doc2_id and doc1_id != doc2_id:
        # Show document info
        col1, col2 = st.columns(2)
        with col1:
            doc1 = st.session_state.documents[doc1_id]
            st.write(f"**Document A:** {doc1['name']}")
            st.write(f"Type: {doc1['type'].upper()}, Words: {doc1['word_count']}")
        
        with col2:
            doc2 = st.session_state.documents[doc2_id]
            st.write(f"**Document B:** {doc2['name']}")
            st.write(f"Type: {doc2['type'].upper()}, Words: {doc2['word_count']}")

        if st.button("Compare Documents", type="primary"):
            if doc1['has_error'] or doc2['has_error']:
                st.error("Cannot compare documents with processing errors.")
                return

            import re
            from collections import Counter
            import difflib

            text1, text2 = doc1['text'], doc2['text']

            # Advanced Text Analysis
            def extract_features(text):
                # Word-level features
                words = re.findall(r'\b\w+\b', text.lower())
                word_freq = Counter(words)
                
                # N-gram analysis
                bigrams = [' '.join(words[i:i+2]) for i in range(len(words)-1)]
                trigrams = [' '.join(words[i:i+3]) for i in range(len(words)-2)]
                
                # Linguistic features
                sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip()]
                avg_sentence_length = sum(len(s.split()) for s in sentences) / len(sentences) if sentences else 0
                
                # Punctuation and formatting
                punctuation = Counter(re.findall(r'[.!?,;:()"\'-]', text))
                
                # Capitalization patterns
                capitals = len(re.findall(r'[A-Z]', text))
                capital_words = len(re.findall(r'\b[A-Z][a-z]*\b', text))
                
                # Numbers and special patterns
                numbers = re.findall(r'\b\d+\b', text)
                emails = re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)
                urls = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', text)
                
                return {
                    'words': set(words), 'word_freq': word_freq, 'bigrams': set(bigrams),
                    'trigrams': set(trigrams), 'sentences': sentences, 'avg_sentence_length': avg_sentence_length,
                    'punctuation': punctuation, 'capitals': capitals, 'capital_words': capital_words,
                    'numbers': numbers, 'emails': emails, 'urls': urls
                }

            features1 = extract_features(text1)
            features2 = extract_features(text2)

            # Semantic similarity using word overlap
            common_words = features1['words'] & features2['words']
            unique_words1 = features1['words'] - features2['words']
            unique_words2 = features2['words'] - features1['words']
            word_similarity = len(common_words) / len(features1['words'] | features2['words']) * 100 if (features1['words'] | features2['words']) else 0

            # Phrase-level analysis
            common_bigrams = set(features1['bigrams']) & set(features2['bigrams'])
            common_trigrams = set(features1['trigrams']) & set(features2['trigrams'])
            
            # Content structure differences
            lines1 = [line.strip() for line in text1.split('\n') if line.strip()]
            lines2 = [line.strip() for line in text2.split('\n') if line.strip()]
            
            # Document sections (assuming headers start with capital letters or numbers)
            sections1 = [line for line in lines1 if re.match(r'^[A-Z0-9].*:?$', line) or len(line.split()) <= 5]
            sections2 = [line for line in lines2 if re.match(r'^[A-Z0-9].*:?$', line) or len(line.split()) <= 5]

            # Display Results
            st.subheader("📊 Advanced Comparison Analysis")
            
            # Overall Metrics
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Word Similarity", f"{word_similarity:.1f}%")
            col2.metric("Phrase Overlap", f"{len(common_bigrams)}")
            col3.metric("Style Similarity", f"{abs(features1['avg_sentence_length'] - features2['avg_sentence_length']):.1f} words diff")
            col4.metric("Structure Diff", f"{abs(len(sections1) - len(sections2))} sections")

            # Lexical Analysis
            st.subheader("🔤 Lexical Analysis")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.write("**Common Elements**")
                st.metric("Shared Words", len(common_words))
                st.metric("Common Phrases", len(common_bigrams))
                st.metric("Common 3-grams", len(common_trigrams))
            
            with col2:
                st.write("**Document A Unique**")
                st.metric("Unique Words", len(unique_words1))
                st.metric("Unique Sentences", len(set(features1['sentences']) - set(features2['sentences'])))
                st.metric("Special Numbers", len(set(features1['numbers']) - set(features2['numbers'])))
            
            with col3:
                st.write("**Document B Unique**")
                st.metric("Unique Words", len(unique_words2))
                st.metric("Unique Sentences", len(set(features2['sentences']) - set(features1['sentences'])))
                st.metric("Special Numbers", len(set(features2['numbers']) - set(features1['numbers'])))

            # Content Type Analysis
            st.subheader("📄 Content Type Differences")
            col1, col2 = st.columns(2)
            
            with col1:
                st.write("**Document A Characteristics**")
                st.write(f"• Average sentence length: {features1['avg_sentence_length']:.1f} words")
                st.write(f"• Capital letters: {features1['capitals']}")
                st.write(f"• Numbers found: {len(features1['numbers'])}")
                st.write(f"• Email addresses: {len(features1['emails'])}")
                st.write(f"• URLs: {len(features1['urls'])}")
                
                if features1['punctuation']:
                    most_common_punct = features1['punctuation'].most_common(3)
                    st.write(f"• Most used punctuation: {', '.join([f'{p}({c})' for p, c in most_common_punct])}")
            
            with col2:
                st.write("**Document B Characteristics**")
                st.write(f"• Average sentence length: {features2['avg_sentence_length']:.1f} words")
                st.write(f"• Capital letters: {features2['capitals']}")
                st.write(f"• Numbers found: {len(features2['numbers'])}")
                st.write(f"• Email addresses: {len(features2['emails'])}")
                st.write(f"• URLs: {len(features2['urls'])}")
                
                if features2['punctuation']:
                    most_common_punct = features2['punctuation'].most_common(3)
                    st.write(f"• Most used punctuation: {', '.join([f'{p}({c})' for p, c in most_common_punct])}")

            # Vocabulary Richness
            st.subheader("📚 Vocabulary Analysis")
            vocab_richness1 = len(set(features1['word_freq'].keys())) / sum(features1['word_freq'].values()) if features1['word_freq'] else 0
            vocab_richness2 = len(set(features2['word_freq'].keys())) / sum(features2['word_freq'].values()) if features2['word_freq'] else 0
            
            col1, col2, col3 = st.columns(3)
            col1.metric("Doc A Vocabulary Richness", f"{vocab_richness1:.3f}")
            col2.metric("Doc B Vocabulary Richness", f"{vocab_richness2:.3f}")
            col3.metric("Richness Difference", f"{abs(vocab_richness1 - vocab_richness2):.3f}")

            # Most frequent unique words
            if unique_words1:
                unique_freq1 = {word: features1['word_freq'][word] for word in unique_words1 if word in features1['word_freq']}
                if unique_freq1:
                    top_unique1 = sorted(unique_freq1.items(), key=lambda x: x[1], reverse=True)[:10]
                    st.write("**Most frequent words unique to Document A:**")
                    st.write(", ".join([f"{word}({freq})" for word, freq in top_unique1]))

            if unique_words2:
                unique_freq2 = {word: features2['word_freq'][word] for word in unique_words2 if word in features2['word_freq']}
                if unique_freq2:
                    top_unique2 = sorted(unique_freq2.items(), key=lambda x: x[1], reverse=True)[:10]
                    st.write("**Most frequent words unique to Document B:**")
                    st.write(", ".join([f"{word}({freq})" for word, freq in top_unique2]))

            # Content Sections
            if sections1 or sections2:
                st.subheader("🏗️ Document Structure")
                col1, col2 = st.columns(2)
                
                with col1:
                    if sections1:
                        st.write("**Document A Sections/Headers:**")
                        for section in sections1[:10]:
                            st.write(f"• {section}")
                
                with col2:
                    if sections2:
                        st.write("**Document B Sections/Headers:**")
                        for section in sections2[:10]:
                            st.write(f"• {section}")

            # Line-by-line differences (for similar documents)
            if word_similarity > 30:
                st.subheader("🔍 Detailed Text Differences")
                
                # Show actual textual differences
                differ = difflib.unified_diff(
                    text1.splitlines(keepends=True)[:20], 
                    text2.splitlines(keepends=True)[:20],
                    fromfile='Document A', 
                    tofile='Document B',
                    lineterm=''
                )
                
                diff_text = ''.join(differ)
                if diff_text:
                    st.code(diff_text, language='diff')
                else:
                    st.info("No significant line-by-line differences found in the first 20 lines.")

            # Overall Assessment
            if word_similarity > 80:
                st.success("🟢 High similarity - documents are very similar!")
            elif word_similarity > 50:
                st.warning("🟡 Moderate similarity - documents share some content")
            else:
                st.info("🔴 Low similarity - documents are quite different")
    else:
        st.warning("Please select two different documents to compare.")

def page_reports():
    st.header("📈 Reports & Insights")
    st.write("Comprehensive analytics across all document formats")

    if not st.session_state.analyses:
        st.info("No analyses performed yet. Run some analyses first.")
        return

    col1, col2, col3, col4 = st.columns(4)

    with col1:
        st.markdown(f"""
            <div style='padding: 1rem; background-color: #fff; border: 1px solid #ccc; border-radius: 10px; text-align: center;'>
                <div style='font-size: 18px; color: #0078d4;'>📄 Total Documents</div>
                <div style='font-size: 32px; font-weight: bold; color: #262730;'>{len(st.session_state.documents)}</div>
            </div>
        """, unsafe_allow_html=True)

    with col2:
        st.markdown(f"""
            <div style='padding: 1rem; background-color: #fff; border: 1px solid #ccc; border-radius: 10px; text-align: center;'>
                <div style='font-size: 18px; color: #0078d4;'>📊 Analyses</div>
                <div style='font-size: 32px; font-weight: bold; color: #262730;'>{len(st.session_state.analyses)}</div>
            </div>
        """, unsafe_allow_html=True)

    with col3:
        file_types = len(set(doc['type'] for doc in st.session_state.documents.values()))
        st.markdown(f"""
            <div style='padding: 1rem; background-color: #fff; border: 1px solid #ccc; border-radius: 10px; text-align: center;'>
                <div style='font-size: 18px; color: #0078d4;'>📋 File Types</div>
                <div style='font-size: 32px; font-weight: bold; color: #262730;'>{file_types}</div>
            </div>
        """, unsafe_allow_html=True)

    with col4:
        avg_score = np.mean([
            analysis['language']['quality_score']
            for analysis in st.session_state.analyses.values()
        ]) if st.session_state.analyses else 0
        st.markdown(f"""
            <div style='padding: 1rem; background-color: #fff; border: 1px solid #ccc; border-radius: 10px; text-align: center;'>
                <div style='font-size: 18px; color: #0078d4;'>✨ Avg Quality</div>
                <div style='font-size: 32px; font-weight: bold; color: #262730;'>{avg_score:.0f}%</div>
            </div>
        """, unsafe_allow_html=True)

    # File type distribution
    if st.session_state.documents:
        st.subheader("📊 Document Type Distribution")
        doc_types = [doc['type'].upper() for doc in st.session_state.documents.values()]
        type_counts = pd.Series(doc_types).value_counts()
        
        col1, col2 = st.columns([2, 1])
        with col1:
            fig = px.bar(x=type_counts.index, y=type_counts.values, 
             title='Document Type Distribution',
             labels={'x': 'File Type', 'y': 'Count'})
            fig.update_xaxes(tickangle=0)
            st.plotly_chart(fig, use_container_width=True)
        with col2:
            st.write("**File Type Summary:**")
            for file_type, count in type_counts.items():
                percentage = (count / len(doc_types)) * 100
                st.write(f"• {file_type}: {count} files ({percentage:.1f}%)")

    st.subheader("📈 Recent Analyses")
    for analysis_id, analysis in list(st.session_state.analyses.items())[-5:]:
        with st.expander(f"📄 {analysis['doc_name']} ({analysis['doc_type'].upper()}) - {analysis['timestamp'].strftime('%Y-%m-%d %H:%M')}"):
            col1, col2, col3 = st.columns(3)
            col2.metric("Sentiment", analysis['nlp']['sentiment'])
            col3.metric("Word Count", analysis['nlp']['word_count'])
            
            if analysis['language']['suggestions']:
                st.write(f"**Suggestions:** {len(analysis['language']['suggestions'])} improvements identified")

def page_health():
    st.header("💚 Enhanced System Health")
    st.write("Monitor system capabilities and file format support")

    # Overall system status
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Status", "Operational", delta="Healthy")
    col2.metric("Documents", len(st.session_state.documents))
    col3.metric("Memory Usage", "Light")
    col4.metric("Processing", "Ready")

    st.subheader("📄 File Format Support Status")
    file_types = get_supported_file_types()
    
    # Group by availability
    available_types = []
    unavailable_types = []
    
    for file_type, info in file_types.items():
        if info['available']:
            available_types.append((file_type, info))
        else:
            unavailable_types.append((file_type, info))

    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### ✅ Available Formats")
        for file_type, info in available_types:
            extensions = ", ".join([f".{ext}" for ext in info['extensions']])
            st.success(f"**{file_type}** ({extensions})")
    
    with col2:
        st.markdown("### ❌ Unavailable Formats")
        if unavailable_types:
            for file_type, info in unavailable_types:
                extensions = ", ".join([f".{ext}" for ext in info['extensions']])
                st.error(f"**{file_type}** ({extensions})")
                st.code(info['install_command'])
        else:
            st.success("All file formats are available! 🎉")

    # Enhanced AI capabilities status
    st.subheader("🤖 AI Enhancement Status")
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### Core AI Features")
        if SPACY_AVAILABLE:
            st.success("✅ spaCy NLP - Entity Recognition Available")
        else:
            st.error("❌ spaCy NLP - Install: pip install spacy")
            st.code("python -m spacy download en_core_web_sm")
        
        if SENTENCE_TRANSFORMER_AVAILABLE:
            st.success("✅ Sentence Transformers - Semantic Analysis Available")
        else:
            st.error("❌ Sentence Transformers - Install: pip install sentence-transformers")
    
    with col2:
        st.markdown("### Enhanced Modules")
        try:
            from enhanced_contradiction_detection import EnhancedContradictionDetector
            st.success("✅ Enhanced Contradiction Detection Module Ready")
        except ImportError:
            st.error("❌ Enhanced Contradiction Detection Module Missing")
            st.write("Save enhanced_contradiction_detection.py in your project directory")

    # Processing capabilities
    st.subheader("🔧 Processing Capabilities")
    capabilities = {
        "Text Extraction": "✅ Multi-format text extraction",
        "NLP Analysis": "✅ Entity recognition, sentiment analysis",
        "Language Analysis": "✅ Quality scoring, improvement suggestions",
        "Enhanced Contradiction Detection": "✅ AI-powered inconsistency detection" if SPACY_AVAILABLE else "⚠️ Basic pattern-based detection",
        "Semantic Analysis": "✅ Meaning-based contradiction detection" if SENTENCE_TRANSFORMER_AVAILABLE else "❌ Requires sentence-transformers",
        "Document Comparison": "✅ Similarity analysis and comparison",
        "Multi-encoding Support": "✅ UTF-8, Latin-1, CP1252 support",
        "Batch Processing": "✅ Multiple file upload and processing",
        "Error Handling": "✅ Graceful error handling and reporting"
    }

    for capability, status in capabilities.items():
        if "✅" in status:
            st.success(f"**{capability}:** {status}")
        elif "⚠️" in status:
            st.warning(f"**{capability}:** {status}")
        else:
            st.error(f"**{capability}:** {status}")

    # File processing statistics
    if st.session_state.documents:
        st.subheader("📊 Processing Statistics")
        
        total_files = len(st.session_state.documents)
        successful_files = sum(1 for doc in st.session_state.documents.values() if not doc['has_error'])
        failed_files = total_files - successful_files
        
        col1, col2, col3 = st.columns(3)
        col1.metric("Total Processed", total_files)
        col2.metric("Successful", successful_files)
        col3.metric("Failed", failed_files)
        
        if total_files > 0:
            success_rate = (successful_files / total_files) * 100
            st.metric("Success Rate", f"{success_rate:.1f}%")

# =============================================================================
# ROUTING (Updated to use the enhanced contradiction detection)
# =============================================================================

if st.session_state.current_page == "📤 Upload Documents":
    page_upload()
elif st.session_state.current_page == "📚 Document Library":
    page_library()
elif st.session_state.current_page == "🔍 Comprehensive Analysis":
    page_comprehensive_analysis()
elif st.session_state.current_page == "⚖️ Contradiction Detection":  # Updated name
    page_contradiction_detection()
elif st.session_state.current_page == "✨ Language Enhancement":
    page_language_enhancement()
elif st.session_state.current_page == "🔀 Document Comparison":
    page_comparison()
elif st.session_state.current_page == "📈 Reports & Insights":
    page_reports()
elif st.session_state.current_page == "📖 User Documentation":
    page_documentation()    
elif st.session_state.current_page == "💚 System Health":
    page_health()
    
# =============================================================================
# FOOTER
# =============================================================================

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #666;'>CoorSight Enhanced Document Analysis Platform - Multi-Format Support with AI-Powered Contradiction Detection</div>",
    unsafe_allow_html=True
)
